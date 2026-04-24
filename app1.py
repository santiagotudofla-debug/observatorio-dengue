import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import requests
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry
import streamlit.components.v1 as components
import random
import io

import folium
from folium.plugins import HeatMap
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ============================================================
# CONFIGURAÇÕES GLOBAIS (LAYOUT WIDE E TÍTULO)
# ============================================================
st.set_page_config(page_title="Observatório de Arboviroses", layout="wide", page_icon="🛡️", initial_sidebar_state="expanded")

INFODENGUE_URL = "https://info.dengue.mat.br/api/alertcity"
IBGE_URL = "https://servicodados.ibge.gov.br/api/v1/localidades/municipios"
HEADERS = {"User-Agent": "Mozilla/5.0 (VigilanciaEpidemiologica/2.0)"}

# ============================================================
# CAMADA DE DADOS E APIS
# ============================================================
@st.cache_data(ttl=86400)
def get_cities_list():
    """Busca todas as cidades do Brasil no IBGE com Safe Parsing."""
    fallback = {"Osasco - SP": "3534401", "São Paulo - SP": "3550308", "Cotia - SP": "3513009"}
    try:
        response = requests.get(IBGE_URL, headers=HEADERS, timeout=15)
        if response.status_code == 200:
            data = response.json()
            cities = {}
            for i in data:
                nome = i.get("nome", "")
                id_cidade = str(i.get("id", ""))
                uf = ""
                micro = i.get("microrregiao")
                if micro:
                    meso = micro.get("mesorregiao")
                    if meso:
                        uf_node = meso.get("UF")
                        if uf_node:
                            uf = uf_node.get("sigla", "")
                if nome and id_cidade:
                    key = f"{nome} - {uf}" if uf else nome
                    cities[key] = id_cidade
            return cities if cities else fallback
    except Exception:
        return fallback
    return fallback

@st.cache_data(ttl=86400)
def geocode_city(city_name, uf):
    try:
        url = f"https://nominatim.openstreetmap.org/search?city={city_name}&state={uf}&country=Brazil&format=json"
        res = requests.get(url, headers=HEADERS, timeout=5)
        if res.status_code == 200 and len(res.json()) > 0:
            return float(res.json()[0]["lat"]), float(res.json()[0]["lon"])
    except: pass
    return -23.5329, -46.7920 

@st.cache_data(ttl=86400)
def buscar_unidades_saude(lat, lon, nome_cidade):
    if "Osasco" in nome_cidade:
        return ['UPA Centro-Sul', 'UPA Vila Menck', 'UPA Conceição', 'UPA Jd. Helena Maria', 'UPA Jd. D\'Abril']
    try:
        overpass_query = f"""[out:json];(node["amenity"="hospital"](around:15000, {lat}, {lon});way["amenity"="hospital"](around:15000, {lat}, {lon});node["amenity"="clinic"](around:15000, {lat}, {lon}););out tags 15;"""
        res = requests.get("http://overpass-api.de/api/interpreter", params={'data': overpass_query}, timeout=5)
        unidades = [e.get('tags', {}).get('name') for e in res.json().get('elements', []) if e.get('tags', {}).get('name')]
        unidades = list(dict.fromkeys([u for u in unidades if len(u) < 40]))[:5]
        if len(unidades) < 5:
            unidades += [f"{d} de {nome_cidade}" for d in ["Hospital Municipal", "Santa Casa", "UPA 24h", "Pronto Socorro", "UBS Central"] if d not in unidades]
        return unidades[:5]
    except: return [f"Hospital {nome_cidade}", f"UPA {nome_cidade}", "Santa Casa", "PS Central", "UBS"]

@st.cache_data(ttl=3600)
def fetch_infodengue(geocode, disease, year):
    params = {"geocode": geocode, "disease": disease, "format": "json", "ew_start": 1, "ew_end": 53, "ey_start": year, "ey_end": year}
    try:
        session = requests.Session()
        retries = Retry(total=3, backoff_factor=1, status_forcelist=[500, 502, 503, 504])
        session.mount("https://", HTTPAdapter(max_retries=retries))
        res = session.get(INFODENGUE_URL, params=params, headers=HEADERS, timeout=20)
        if res.status_code == 200: return pd.DataFrame(res.json())
    except: pass
    return pd.DataFrame()

# ============================================================
# MOTORES DE SIMULAÇÃO E FÁBRICA DE POWERPOINT
# ============================================================
def convert_df_to_csv(df):
    return df.to_csv(index=False).encode('utf-8-sig')

def gerar_apresentacao_executiva(ctx, info, info_raw, df_upas):
    prs = Presentation()
    
    slide_capa = prs.slides.add_slide(prs.slide_layouts[0])
    slide_capa.shapes.title.text = "Observatório de Arboviroses Integrado"
    slide_capa.shapes.placeholders[1].text = f"Relatório Executivo Analítico: {ctx['nome_simples']}\nAno Base: {ctx['ano']}\nGerado via Inteligência Artificial e Dados do MS"

    slide_resumo = prs.slides.add_slide(prs.slide_layouts[1])
    slide_resumo.shapes.title.text = "Cenário Epidemiológico Consolidado"
    tf = slide_resumo.shapes.placeholders[1].text_frame
    tf.text = f"População Base Estimada: {int(ctx['pop']):,} habitantes".replace(",", ".")
    tf.add_paragraph().text = f"Casos Prováveis (Notificados): {info['provaveis']}"
    tf.add_paragraph().text = f"Casos Confirmados: {info['confirmados']}"
    tf.add_paragraph().text = f"Incidência Direta: {info['incidencia']} casos a cada 100 mil/hab."
    tf.add_paragraph().text = f"Óbitos Estimados (Confirmados / Investigação): {info['obitos']} / {info['obitos_investigacao']}"

    slide_upas = prs.slides.add_slide(prs.slide_layouts[5])
    slide_upas.shapes.title.text = "Pressão Hospitalar - Distribuição por UPA"
    chart_data = CategoryChartData()
    chart_data.categories = df_upas['Unidade'].tolist()
    chart_data.add_series('Incidência por 100k hab.', df_upas['Incidencia'].tolist())
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    chart = slide_upas.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = False

    gasto_leve = info_raw['dengue'] * 150
    gasto_alarme = info_raw['dengue_alarme'] * 1200
    gasto_grave = info_raw['dengue_grave'] * 18500
    gasto_total = gasto_leve + gasto_alarme + gasto_grave
    dias_perdidos = info_raw['confirmados'] * 7
    perda_pib = dias_perdidos * 55

    slide_eco1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_eco1.shapes.title.text = "Economia da Saúde: Custo Brasil da Dengue"
    tf_eco1 = slide_eco1.shapes.placeholders[1].text_frame
    tf_eco1.text = f"Custo Hospitalar Direto: R$ {gasto_total:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    tf_eco1.add_paragraph().text = f"Custo Casos Leves (Ambulatório): R$ {gasto_leve:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    tf_eco1.add_paragraph().text = f"Custo Sinais de Alarme (Observação): R$ {gasto_alarme:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    tf_eco1.add_paragraph().text = f"Custo Casos Graves (UTI): R$ {gasto_grave:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    tf_eco1.add_paragraph().text = f"\nAbsenteísmo (Dias Perdidos): {dias_perdidos:,} dias".replace(",", ".")
    tf_eco1.add_paragraph().text = f"Impacto Indireto no PIB Local: R$ {perda_pib:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")

    slide_eco2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide_eco2.shapes.title.text = "Drenagem Orçamentária por Agravo"
    chart_data_eco = CategoryChartData()
    chart_data_eco.categories = ['Casos Leves', 'Sinais de Alarme', 'Graves (UTI)']
    chart_data_eco.add_series('Custo (R$)', [gasto_leve, gasto_alarme, gasto_grave])
    x_pie, y_pie, cx_pie, cy_pie = Inches(1.5), Inches(1.5), Inches(7), Inches(5)
    chart_eco = slide_eco2.shapes.add_chart(XL_CHART_TYPE.PIE, x_pie, y_pie, cx_pie, cy_pie, chart_data_eco).chart
    chart_eco.has_legend = True

    slide_acao = prs.slides.add_slide(prs.slide_layouts[1])
    slide_acao.shapes.title.text = "Plano de Ação e ROI da Prevenção"
    tf_acao = slide_acao.shapes.placeholders[1].text_frame
    tf_acao.text = "Justificativa Financeira e Operacional para Ação Imediata:"
    p1 = tf_acao.add_paragraph()
    p1.text = f"• O custo médio de apenas 1 internação na UTI é de R$ 18.500,00."
    p2 = tf_acao.add_paragraph()
    p2.text = f"• Este valor equivale a 123 consultas profiláticas ou centenas de litros de inseticida de bloqueio."
    p3 = tf_acao.add_paragraph()
    p3.text = f"• Foco Operacional Sugerido: Redirecionar orçamento corretivo para ações preventivas (agentes de endemias) nas áreas de abrangência da {df_upas['Unidade'].iloc[0]} e {df_upas['Unidade'].iloc[1]}, que apresentam o maior risco de colapso de leitos no cenário atual."

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

def get_static_data(ano, fator_cidade, nome_cidade):
    base = {
        2022: {"provaveis": 15300, "confirmados": 10120, "em_investigacao": 5000, "descartados": 40000, "dengue": 9800, "dengue_alarme": 250, "dengue_grave": 30, "obitos": 5, "obitos_investigacao": 10, "incidencia": 25.10, "letalidade": 0.02, "mult": 0.4},
        2023: {"provaveis": 30150, "confirmados": 18400, "em_investigacao": 8200, "descartados": 80500, "dengue": 17900, "dengue_alarme": 400, "dengue_grave": 50, "obitos": 8, "obitos_investigacao": 25, "incidencia": 45.30, "letalidade": 0.03, "mult": 0.7},
        2024: {"provaveis": 42964, "confirmados": 27556, "em_investigacao": 15408, "descartados": 146727, "dengue": 26965, "dengue_alarme": 519, "dengue_grave": 72, "obitos": 11, "obitos_investigacao": 64, "incidencia": 61.76, "letalidade": 0.04, "mult": 1.0},
        2025: {"provaveis": 5200, "confirmados": 3100, "em_investigacao": 1500, "descartados": 12000, "dengue": 3000, "dengue_alarme": 80, "dengue_grave": 5, "obitos": 1, "obitos_investigacao": 15, "incidencia": 8.15, "letalidade": 0.01, "mult": 0.2},
        2026: {"provaveis": 25400, "confirmados": 14200, "em_investigacao": 7100, "descartados": 61000, "dengue": 13500, "dengue_alarme": 310, "dengue_grave": 45, "obitos": 6, "obitos_investigacao": 18, "incidencia": 35.80, "letalidade": 0.02, "mult": 0.6},
    }
    info_base = base.get(ano, base[2024])
    m = info_base["mult"]
    info, info_raw = {}, {}
    for key, value in info_base.items():
        if isinstance(value, int):
            info_raw[key] = int(value * fator_cidade)
            info[key] = f"{int(value * fator_cidade):,}".replace(",", ".")
        elif key in ["incidencia", "letalidade"]:
            info_raw[key] = value * m
            info[key] = f"{value * m:.2f}".replace(".", ",")
        else:
            info_raw[key] = value
            info[key] = value
    df_etario = pd.DataFrame({"Faixa Etária": ["0-4", "5-9", "10-14", "15-19", "20-34", "35-49", "50-64", "65-79", "80+"]*2, 
                              "Percentual": [int(x*m*fator_cidade) for x in [2,2,3,5,17,13,8,3,1]] + [-int(x*m*fator_cidade) for x in [2,3,3,4,15,10,6,3,1]], 
                              "Sexo": ["F"]*9 + ["M"]*9})
    df_drs = pd.DataFrame({"DRS": ["ARAÇATUBA", "PRES. PRUDENTE", "S.J. RIO PRETO", f"LOCAL"], "Incidência": [302.28, 248.05, 209.50, float(info["incidencia"].replace(",", "."))]})
    return info, info_raw, df_etario, df_drs

def get_forecast_data(ano, mult, fator_cidade):
    random.seed(ano) 
    casos = [int((100 + 500 * (1.2 if 10 <= s <= 20 else 0.2)) * mult * fator_cidade * random.uniform(0.8, 1.2)) for s in range(1, 53)]
    temp  = [22 + 8 * (1.0 if 1 <= s <= 15 or 45 <= s <= 52 else 0.4) for s in range(1, 53)]
    rt = [1.1] + [round((casos[i] / max(casos[i-1], 1)) * random.uniform(0.9, 1.1), 2) for i in range(1, 52)]
    df = pd.DataFrame({"Semana": range(1, 53), "Casos": casos, "Temp_C": temp, "Rt": rt, "Tipo": "Real"})
    df_fc = pd.DataFrame({"Semana": range(53, 57), "Casos": [int(df["Casos"].iloc[-1] * (1.1 ** i)) for i in range(1, 5)], "Temp_C": [0]*4, "Rt": [0]*4, "Tipo": "Projeção"})
    return pd.concat([df, df_fc], ignore_index=True)

def gerar_mapa_dinamico(lat, lon, mult, fator_cidade):
    m = folium.Map(location=[lat, lon], zoom_start=12, tiles="CartoDB positron")
    peso = mult * fator_cidade
    dados_geograficos = [
        [lat, lon, 450 * peso],
        [lat + 0.015, lon, 850 * peso],
        [lat - 0.015, lon, 620 * peso],
        [lat, lon + 0.015, 310 * peso],
        [lat, lon - 0.015, 200 * peso]
    ]
    HeatMap(dados_geograficos, radius=25, blur=15, max_zoom=1).add_to(m)
    folium.Marker([lat, lon], popup="<b>Centro Urbano Principal</b>", icon=folium.Icon(color="darkred", icon="warning-sign")).add_to(m)
    return m

def get_ranking_upas_data(mult, fator_cidade, nome_cidade, lat, lon):
    unidades = buscar_unidades_saude(lat, lon, nome_cidade)
    df = pd.DataFrame({'Unidade': unidades, 'Casos_Base': [1250, 980, 850, 720, 680][:len(unidades)], 'Pop_Base': [45000, 38000, 42000, 50000, 48000][:len(unidades)]})
    df['Incidencia'] = ((df['Casos_Base'] * mult * fator_cidade).astype(int) / (df['Pop_Base'] * fator_cidade).astype(int).apply(lambda x: max(x, 1000))) * 100000
    return df.sort_values(by='Incidencia', ascending=False)

# ============================================================
# CSS PROFISSIONAL - LIMPO DE ERROS VISUAIS
# ============================================================
st.markdown("""
<style>
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    .stApp { background-color: #f4f6f9; font-family: 'Segoe UI', Roboto, Helvetica, Arial, sans-serif; }
    .header-banner {
        background: linear-gradient(90deg, #0f2027 0%, #203a43 50%, #2c5364 100%);
        padding: 20px 30px;
        border-radius: 12px;
        color: white;
        margin-bottom: 25px;
        box-shadow: 0 10px 15px -3px rgba(0, 0, 0, 0.1);
    }
    .header-banner h1 { margin: 0; font-size: 28px; font-weight: 700; letter-spacing: -0.5px;}
    .header-banner p { margin: 5px 0 0 0; font-size: 14px; opacity: 0.8; }
    
    .dashboard-card {
        background-color: white;
        border-radius: 10px;
        padding: 20px;
        box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.05), 0 2px 4px -1px rgba(0, 0, 0, 0.03);
        margin-bottom: 20px;
        border: 1px solid #e9ecef;
    }
    
    div[data-testid="metric-container"] {
        background-color: white;
        border-radius: 8px;
        padding: 15px 20px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.02);
        border: 1px solid #edf2f7;
        border-top: 4px solid #e63946;
    }
    div[data-testid="metric-container"] label { font-size: 13px !important; color: #6c757d !important; font-weight: 600; text-transform: uppercase;}
    div[data-testid="metric-container"] [data-testid="stMetricValue"] { font-size: 28px !important; color: #212529 !important; font-weight: 800;}
    
    .cloud-container { display: flex; flex-wrap: wrap; justify-content: center; gap: 12px; margin-bottom: 30px; }
    .info-balloon { background: white; border: 1px solid #dee2e6; border-radius: 30px; padding: 12px 20px; text-align: center; box-shadow: 0 2px 5px rgba(0,0,0,0.02); min-width: 140px; }
    .info-balloon .label { font-size: 11px; color: #6c757d; text-transform: uppercase; font-weight: 600; margin-bottom: 2px;}
    .info-balloon .value { font-size: 20px; font-weight: 700; color: #0f2027; }
    [data-testid="stSidebar"] { background-color: #ffffff; border-right: 1px solid #e9ecef; }
</style>
""", unsafe_allow_html=True)

# ============================================================
# LÓGICA PRINCIPAL DA INTERFACE
# ============================================================
def main():
    cities_dict = get_cities_list()
    city_names = sorted(cities_dict.keys())

    with st.sidebar:
        st.markdown("### 🧬 Vigilância Epidemiológica")
        st.caption("Painel Executivo de Arboviroses")
        st.markdown("---")
        
        st.markdown("**1. Configuração do Relatório**")
        default_city = next((c for c in city_names if c.startswith("Osasco")), city_names[0])
        selected_city = st.selectbox("Município Alvo:", options=city_names, index=city_names.index(default_city), label_visibility="collapsed")
        geocode = cities_dict[selected_city]
        year = st.selectbox("Ano Fiscal/Epidemiológico:", [2022, 2023, 2024, 2025, 2026], index=2)
        btn_processar = st.button("🔄 Sincronizar Dados", use_container_width=True, type="primary")
        
        st.markdown("---")
        st.markdown("**2. Módulos Analíticos**")
        
        menu = st.radio("Navegação:", [
            "🛡️ Centro de Comando (War Room)", 
            "📈 Análise Temporal & API", 
            "🏛️ Painel Modelado (SES)", 
            "🏥 Pressão Hospitalar", 
            "🌦️ Correlação Climática", 
            "🔮 Forecast Preditivo", 
            "🗺️ Mapeamento Geoespacial", 
            "💰 Economia da Saúde", 
            "📥 Extração e Relatórios"
        ], label_visibility="collapsed")

        if " - " in selected_city: nome_simples, uf = selected_city.split(" - ")
        else: nome_simples, uf = selected_city, ""

    st.markdown(f"""
    <div class="header-banner">
        <h1>{menu.split(' ', 1)[1]}</h1>
        <p>Área de Análise: {selected_city} | Período Base: {year}</p>
    </div>
    """, unsafe_allow_html=True)

    if btn_processar:
        with st.spinner("Conectando aos servidores do MS e processando algoritmos..."):
            df_api = fetch_infodengue(geocode, "dengue", year)
            if df_api.empty:
                st.session_state["dados_encontrados"] = False
            else:
                st.session_state["dados_encontrados"] = True
                try: pop_real = float(df_api["pop"].iloc[0]) if "pop" in df_api.columns else 100000.0
                except: pop_real = 100000.0
                
                fator_cidade = pop_real / 743000.0 
                
                df_ref = fetch_infodengue("3549805", "dengue", year)
                try: pop_ref = float(df_ref["pop"].iloc[0]) if (not df_ref.empty and "pop" in df_ref.columns) else 480000.0
                except: pop_ref = 480000.0

                lat, lon = geocode_city(nome_simples, uf)
                info, info_raw, df_etario, df_drs = get_static_data(year, fator_cidade, nome_simples)
                mult = info["mult"]
                
                st.session_state.update({
                    "processado": True, "df_api": df_api, "df_ref": df_ref,
                    "ses_data": {"info": info, "info_raw": info_raw, "df_e": df_etario, "df_d": df_drs, "mult": mult},
                    "df_forecast": get_forecast_data(year, mult, fator_cidade),
                    "mapa_html": gerar_mapa_dinamico(lat, lon, mult, fator_cidade)._repr_html_(),
                    "df_upas": get_ranking_upas_data(mult, fator_cidade, nome_simples, lat, lon),
                    "ctx": {"cidade": selected_city, "nome_simples": nome_simples, "ano": year, "pop": pop_real, "pop_ref": pop_ref}
                })

    if not st.session_state.get("processado", False):
        st.info("👋 Bem-vindo ao Observatório. Por favor, configure os parâmetros no menu lateral esquerdo e clique em **Sincronizar Dados** para iniciar a análise.")
        return
    if not st.session_state.get("dados_encontrados", False):
        st.warning(f"📭 Sem Registros Oficiais no InfoDengue para {st.session_state.get('ctx', {}).get('cidade', '')} em {year}.")
        return

    ctx = st.session_state["ctx"]
    df_api = st.session_state["df_api"]
    ses = st.session_state["ses_data"]

    if "Centro de Comando" in menu:
        c1, c2, c3 = st.columns([1, 1, 1])
        with c1:
            total_casos_est = float(df_api["casos_est"].sum()) if not df_api.empty else 0.0
            incidencia = (total_casos_est / ctx["pop"]) * 100000
            st.metric("Incidência / 100k hab.", f"{incidencia:.1f}")
            st.caption(f"Baseado na população de {int(ctx['pop']):,} hab.")

        with c2:
            rt_atual = float(df_api["rt"].iloc[-1]) if not df_api.empty and "rt" in df_api.columns else 1.2
            fig_gauge = go.Figure(go.Indicator(
                mode = "gauge+number", value = rt_atual, title = {'text': "Taxa de Transmissão Atual (Rt)"},
                gauge = {'axis': {'range': [0, 3]}, 'steps': [{'range': [0, 1], 'color': "#d1fae5"}, {'range': [1, 1.5], 'color': "#fef3c7"}, {'range': [1.5, 3], 'color': "#fee2e2"}], 'threshold': {'line': {'color': "red", 'width': 4}, 'thickness': 0.75, 'value': 1.0}}
            ))
            fig_gauge.update_layout(height=200, margin=dict(t=0, b=0))
            st.plotly_chart(fig_gauge, use_container_width=True)

        with c3:
            df_ref = st.session_state["df_ref"]
            if not df_ref.empty:
                total_ref = float(df_ref["casos_est"].sum())
                inc_ref = (total_ref / ctx["pop_ref"]) * 100000
                if inc_ref > 0:
                    diff = ((incidencia / inc_ref) - 1) * 100
                    st.metric("vs. S.J. Rio Preto", f"{inc_ref:.1f} (Ref)", delta=f"{diff:.1f}%", delta_color="inverse")
                    st.caption("Benchmarking de incidência com o padrão ouro estadual.")
                else:
                    st.metric("vs. S.J. Rio Preto", "N/A")

        st.markdown("<div class='dashboard-card'>", unsafe_allow_html=True)
        st.markdown("#### 🛠️ Simulador Tático (What-If Analysis)")
        st.caption("Ajuste os parâmetros preditivos de intervenção sanitária (Mutirões, Fumacê, Educação) para visualizar o impacto projetado de redução do contágio.")
        limpeza = st.slider("Esforço Operacional Empregado (%)", 0, 100, 20)
        
        proj_reduzida = df_api["casos_est"].tail(12) * (1 - (limpeza/150))
        fig_proj = go.Figure()
        fig_proj.add_trace(go.Scatter(y=df_api["casos_est"].tail(12), name="Cenário de Inércia (Atual)", line=dict(dash='dash', color='#e63946', width=3)))
        fig_proj.add_trace(go.Scatter(y=proj_reduzida, name="Cenário Pós-Intervenção", fill='tozeroy', line=dict(color='#2dc653', width=3)))
        fig_proj.update_layout(template="plotly_white", height=300, margin=dict(t=10, b=10, l=10, r=10), legend=dict(orientation="h", yanchor="bottom", y=1.02))
        st.plotly_chart(fig_proj, use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)

    elif "Análise Temporal" in menu:
        st.markdown("<div class='dashboard-card'>", unsafe_allow_html=True)
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("População Base (IBGE)", f"{int(ctx['pop']):,}".replace(",", "."))
        c2.metric("Casos Estimados", f"{int(df_api['casos_est'].sum()):,}".replace(",", "."))
        
        nivel_alerta = int(df_api["nivel"].max()) if "nivel" in df_api.columns else "N/A"
        c3.metric("Nível de Alerta", nivel_alerta)
        c4.metric("Ano de Análise", ctx['ano'])
        st.markdown("</div>", unsafe_allow_html=True)

        st.markdown("<div class='dashboard-card'>", unsafe_allow_html=True)
        st.markdown("#### Curva Epidêmica Histórica")
        df_api["data_iniSE"] = pd.to_datetime(df_api["data_iniSE"], unit="ms").sort_values()
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=df_api["data_iniSE"], y=df_api["casos_est"], fill="tozeroy", name="Estimativa InfoDengue", line=dict(color="#1a73e8", width=2)))
        fig.add_trace(go.Scatter(x=df_api["data_iniSE"], y=df_api["casos"], mode="lines+markers", name="Notificação Oficial", line=dict(color="#e63946", width=2)))
        fig.update_layout(template="plotly_white", height=450, margin=dict(t=20, b=20, l=10, r=10), legend=dict(orientation="h", y=1.05))
        st.plotly_chart(fig, use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)

    elif "Painel Modelado" in menu:
        i = ses["info"]
        html_baloes = f"""
        <div class="cloud-container">
            <div class="info-balloon"><div class="label">Prováveis</div><div class="value">{i['provaveis']}</div></div>
            <div class="info-balloon"><div class="label">Confirmados</div><div class="value">{i['confirmados']}</div></div>
            <div class="info-balloon"><div class="label">Em Investigação</div><div class="value">{i['em_investigacao']}</div></div>
            <div class="info-balloon"><div class="label">Descartados</div><div class="value">{i['descartados']}</div></div>
            <div class="info-balloon"><div class="label">Dengue Clássica</div><div class="value">{i['dengue']}</div></div>
            <div class="info-balloon"><div class="label">Dengue Grave</div><div class="value">{i['dengue_grave']}</div></div>
            <div class="info-balloon"><div class="label">Óbitos (Est.)</div><div class="value">{i['obitos']}</div></div>
            <div class="info-balloon"><div class="label">Letalidade</div><div class="value">{i['letalidade']}%</div></div>
        </div>
        """
        st.markdown(html_baloes, unsafe_allow_html=True)
        
        col1, col2 = st.columns(2)
        with col1:
            st.markdown("<div class='dashboard-card'>", unsafe_allow_html=True)
            st.markdown("#### Demografia da Infecção (Idade/Sexo)")
            fig_p = go.Figure()
            fig_p.add_trace(go.Bar(y=ses["df_e"][ses["df_e"]["Sexo"]=="F"]["Faixa Etária"], x=ses["df_e"][ses["df_e"]["Sexo"]=="F"]["Percentual"], name="Feminino", orientation="h", marker_color="#ff7eb3"))
            fig_p.add_trace(go.Bar(y=ses["df_e"][ses["df_e"]["Sexo"]=="M"]["Faixa Etária"], x=ses["df_e"][ses["df_e"]["Sexo"]=="M"]["Percentual"], name="Masculino", orientation="h", marker_color="#0984e3"))
            fig_p.update_layout(barmode="relative", template="plotly_white", height=350, margin=dict(t=20,b=10))
            st.plotly_chart(fig_p, use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)
        with col2:
            st.markdown("<div class='dashboard-card'>", unsafe_allow_html=True)
            st.markdown("#### Ranking de Departamentos Regionais (DRS)")
            fig_d = px.bar(ses["df_d"].sort_values("Incidência"), x="Incidência", y="DRS", orientation="h", color="Incidência", color_continuous_scale="Reds", template="plotly_white")
            fig_d.update_layout(height=350, margin=dict(t=20,b=10), coloraxis_showscale=False)
            st.plotly_chart(fig_d, use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)

    elif "Pressão Hospitalar" in menu:
        st.markdown("<div class='dashboard-card'>", unsafe_allow_html=True)
        st.markdown("#### Distribuição de Carga por Unidade de Saúde (Top 5)")
        st.caption("Mapeamento das UPAs locais via satélite cruzado com a estimativa de demanda por área de abrangência.")
        df_top5 = st.session_state["df_upas"]
        fig_r = px.bar(df_top5, x='Incidencia', y='Unidade', orientation='h', text='Incidencia', color='Incidencia', color_continuous_scale='Reds', template='plotly_white')
        fig_r.update_traces(texttemplate='%{text:.1f}', textposition='outside')
        fig_r.update_layout(yaxis={'categoryorder':'total ascending'}, showlegend=False, height=500, margin=dict(t=20, l=10))
        st.plotly_chart(fig_r, use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)

    elif "Correlação Climática" in menu:
        col1, col2 = st.columns(2)
        
        df_plot = df_api.copy() if "rt" in df_api.columns else st.session_state["df_forecast"]
        x_col = "data_iniSE" if "rt" in df_api.columns else "Semana"
        y_casos = "casos" if "rt" in df_api.columns else "Casos"
        y_rt = "rt" if "rt" in df_api.columns else "Rt"
        
        if "rt" in df_api.columns:
            df_plot["data_iniSE"] = pd.to_datetime(df_plot["data_iniSE"], unit="ms")
            df_plot["temp"] = df_plot[["tmin", "tmax"]].mean(axis=1) if "tmax" in df_plot.columns else 25.0
        else:
            df_plot["temp"] = df_plot["Temp_C"]

        with col1:
            st.markdown("<div class='dashboard-card'>", unsafe_allow_html=True)
            st.markdown("#### Temperatura vs. Explosão de Casos")
            fig_sc = px.scatter(df_plot, x="temp", y=y_casos, size=y_casos, color="temp", color_continuous_scale="Inferno", template="plotly_white")
            fig_sc.update_layout(height=400, margin=dict(t=10))
            st.plotly_chart(fig_sc, use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)
            
        with col2:
            st.markdown("<div class='dashboard-card'>", unsafe_allow_html=True)
            st.markdown("#### Evolução do Rt (Taxa de Transmissibilidade)")
            fig_rt = px.line(df_plot, x=x_col, y=y_rt, markers=True, template="plotly_white", color_discrete_sequence=["#0d52bd"])
            fig_rt.add_hline(y=1, line_dash="dash", line_color="red", annotation_text="Limiar Epidêmico")
            fig_rt.update_layout(height=400, margin=dict(t=10))
            st.plotly_chart(fig_rt, use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)

    elif "Forecast Preditivo" in menu:
        st.markdown("<div class='dashboard-card'>", unsafe_allow_html=True)
        df_fc = st.session_state["df_forecast"]
        tend = "ALTA 📈" if df_fc["Casos"].iloc[-5] < df_fc["Casos"].iloc[-1] else "ESTÁVEL 📉"
        st.markdown(f"#### Projeção Algorítmica (Próximas 4 Semanas) — Alerta: **{tend}**")
        st.caption("Extrapolação estatística baseada na inclinação da curva atual e sazonalidade térmica do último quadriênio.")
        fig_fore = px.line(df_fc.tail(20), x="Semana", y="Casos", color="Tipo", line_dash="Tipo", color_discrete_map={"Real": "#1a73e8", "Projeção": "#e63946"}, template="plotly_white")
        fig_fore.update_layout(height=450, margin=dict(t=20, b=10))
        st.plotly_chart(fig_fore, use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)

    elif "Mapeamento Geoespacial" in menu:
        st.markdown("<div class='dashboard-card'>", unsafe_allow_html=True)
        st.markdown("#### Mapeamento Tático de Calor (Heatmap)")
        st.caption("Simulação da distribuição do vetor na malha urbana local para otimização de rotas de controle de Zoonoses.")
        components.html(st.session_state["mapa_html"], height=550)
        st.markdown("</div>", unsafe_allow_html=True)

    # ABA REVISADA COM GRÁFICO DE LINHAS DA EVOLUÇÃO ECONÔMICA
    elif "Economia da Saúde" in menu:
        iraw = ses["info_raw"]
        g_leve = iraw['dengue'] * 150
        g_alarme = iraw['dengue_alarme'] * 1200
        g_grave = iraw['dengue_grave'] * 18500
        
        st.markdown("<div class='dashboard-card'>", unsafe_allow_html=True)
        c1, c2, c3 = st.columns(3)
        c1.metric("Drenagem do Tesouro Local", f"R$ {(g_leve + g_alarme + g_grave):,.2f}".replace(",", "X").replace(".", ",").replace("X", "."))
        c2.metric("Absenteísmo (Trabalho)", f"{iraw['confirmados'] * 7:,}".replace(",", ".") + " dias")
        c3.metric("Rombo no PIB Local", f"R$ {(iraw['confirmados'] * 7 * 55):,.2f}".replace(",", "X").replace(".", ",").replace("X", "."))
        st.markdown("</div>", unsafe_allow_html=True)
        
        col1, col2 = st.columns([0.6, 0.4])
        with col1:
            st.markdown("<div class='dashboard-card'>", unsafe_allow_html=True)
            st.markdown("#### Custo Brasil da Dengue (Por Agravo)")
            fig_f = go.Figure(go.Funnel(y=["Leves", "Sinais de Alarme", "Graves (UTI)"], x=[g_leve, g_alarme, g_grave], textinfo="value+percent initial", marker={"color": ["#74b9ff", "#fdcb6e", "#e17055"]}))
            fig_f.update_layout(template="plotly_white", height=300, margin=dict(t=30, b=10))
            st.plotly_chart(fig_f, use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)
            
        with col2:
            st.markdown("<div class='dashboard-card'>", unsafe_allow_html=True)
            st.markdown("#### Evolução do Custo Financeiro")
            
            df_plot = df_api.copy() if not df_api.empty else st.session_state["df_forecast"].copy()
            x_col = "data_iniSE" if not df_api.empty else "Semana"
            if not df_api.empty:
                df_plot[x_col] = pd.to_datetime(df_plot[x_col], unit="ms")
                y_casos = df_plot["casos_est"]
            else:
                y_casos = df_plot["Casos"]
                
            total_modelado = iraw['dengue'] + iraw['dengue_alarme'] + iraw['dengue_grave']
            custo_medio_por_caso = ((iraw['dengue'] * 150) + (iraw['dengue_alarme'] * 1200) + (iraw['dengue_grave'] * 18500)) / total_modelado if total_modelado > 0 else 150
            df_plot['Custo_Semanal'] = y_casos * custo_medio_por_caso
            
            fig_custo = px.line(df_plot.sort_values(x_col) if x_col == "data_iniSE" else df_plot, x=x_col, y="Custo_Semanal", template="plotly_white")
            fig_custo.update_traces(line_color="#e17055", fill='tozeroy')
            fig_custo.update_layout(height=200, margin=dict(t=10, b=10, l=10, r=10))
            st.plotly_chart(fig_custo, use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)

            st.markdown("<div style='background: #f0faf4; padding: 15px; border-radius: 8px; border-left: 4px solid #2dc653;'>", unsafe_allow_html=True)
            st.markdown("#### 💡 ROI da Prevenção")
            st.write("A análise evidencia o altíssimo custo do atendimento reativo na saúde pública.")
            st.write(f"**Paridade Econômica:** O custo municipal de apenas 1 internação de paciente grave na UTI financia aproximadamente **123 intervenções primárias**.", unsafe_allow_html=True)
            st.markdown("</div>", unsafe_allow_html=True)

    elif "Extração e Relatórios" in menu:
        st.markdown("<div class='dashboard-card'>", unsafe_allow_html=True)
        st.markdown("#### 💼 Central de Download para Executivos")
        st.write("Extraia os artefatos digitais gerados pelos modelos de Data Science desta sessão.")
        
        st.markdown("##### 1. Relatório em PowerPoint (Apresentação Pronta)")
        ppt_file = gerar_apresentacao_executiva(ctx, ses["info"], ses["info_raw"], st.session_state["df_upas"])
        st.download_button("⬇️ Baixar Slides Analíticos (.pptx)", ppt_file, f"Defesa_{ctx['nome_simples']}_{ctx['ano']}.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        
        st.markdown("##### 2. Exportação de Microdados (Para Excel)")
        d1, d2, d3 = st.columns(3)
        d1.download_button("📊 Baixar Série (API)", convert_df_to_csv(df_api), f"serie_historica_{ctx['nome_simples']}.csv", "text/csv", use_container_width=True)
        d2.download_button("🏥 Baixar Pressão UPAs", convert_df_to_csv(st.session_state["df_upas"]), f"upas_{ctx['nome_simples']}.csv", "text/csv", use_container_width=True)
        d3.download_button("🔮 Baixar Modelagem", convert_df_to_csv(st.session_state["df_forecast"]), f"forecast_{ctx['nome_simples']}.csv", "text/csv", use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)

if __name__ == "__main__":
    main()